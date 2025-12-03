#!/usr/bin/env python3
"""
Script de génération de la présentation PowerPoint
Orange Think Tank Challenge 2025 - Prédiction de la propagation des épidémies

Ce script crée automatiquement une présentation PowerPoint professionnelle
conforme aux critères d'évaluation du challenge.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

# =============================================================================
# CONFIGURATION - Charte Graphique Orange CI
# =============================================================================

ORANGE_CI = RGBColor(250, 126, 25)      # #fa7e19
NOIR = RGBColor(0, 0, 0)                 # #000000
GRIS_CLAIR = RGBColor(247, 247, 247)     # #F7F7F7
BLANC = RGBColor(255, 255, 255)          # #FFFFFF

# Chemins des images
ARCHITECTURE_IMAGE = "architecture.png"
ANONYMISATION_IMAGE = "anonymisation.png"

# =============================================================================
# FONCTIONS UTILITAIRES
# =============================================================================

def add_title(slide, title_text, font_size=44):
    """Ajoute un titre formaté à une slide"""
    # Créer une textbox pour le titre (car nous utilisons un layout blank)
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(9), Inches(0.8)
    )
    text_frame = title_box.text_frame
    text_frame.word_wrap = True

    p = text_frame.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(font_size)
    p.font.bold = True
    p.font.color.rgb = ORANGE_CI
    p.alignment = PP_ALIGN.CENTER

    return title_box


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=NOIR, alignment=PP_ALIGN.LEFT):
    """Ajoute une zone de texte formatée"""
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.word_wrap = True

    p = text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment

    return textbox


def add_shape_with_text(slide, left, top, width, height, text,
                        fill_color=GRIS_CLAIR, line_color=ORANGE_CI,
                        font_size=16, font_color=NOIR):
    """Ajoute une forme rectangulaire avec texte"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height
    )

    # Formatage de la forme
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = line_color
    shape.line.width = Pt(2)

    # Ajout du texte
    text_frame = shape.text_frame
    text_frame.word_wrap = True
    text_frame.margin_left = Inches(0.1)
    text_frame.margin_right = Inches(0.1)
    text_frame.margin_top = Inches(0.05)

    p = text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.alignment = PP_ALIGN.LEFT

    return shape


def add_page_number(slide, slide_number, total_slides=7):
    """Ajoute le numéro de page en bas à droite"""
    page_text = f"{slide_number}/{total_slides}"
    textbox = slide.shapes.add_textbox(
        Inches(8.5), Inches(7), Inches(1), Inches(0.3)
    )
    p = textbox.text_frame.paragraphs[0]
    p.text = page_text
    p.font.size = Pt(12)
    p.font.color.rgb = NOIR
    p.alignment = PP_ALIGN.RIGHT


def add_bullet_point(text_frame, text, level=0, font_size=16):
    """Ajoute un point de liste"""
    p = text_frame.add_paragraph()
    p.text = text
    p.level = level
    p.font.size = Pt(font_size)
    p.font.color.rgb = NOIR
    return p


# =============================================================================
# CRÉATION DES SLIDES
# =============================================================================

def create_slide_1_title(prs):
    """
    SLIDE 1: Page de Titre + Problématique
    """
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Fond gris clair
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = GRIS_CLAIR

    # Titre principal
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.5), Inches(9), Inches(1.5)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Prédire la Propagation des Épidémies"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = ORANGE_CI
    p.alignment = PP_ALIGN.CENTER

    # Sous-titre
    subtitle_box = slide.shapes.add_textbox(
        Inches(1), Inches(3), Inches(8), Inches(0.7)
    )
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Transformer les données de mobilité en outil de santé publique"
    p.font.size = Pt(28)
    p.font.color.rgb = NOIR
    p.alignment = PP_ALIGN.CENTER

    # Encadré problématique (orange)
    problem_box = add_shape_with_text(
        slide,
        Inches(1.5), Inches(4.5), Inches(7), Inches(1.2),
        "Comment anticiper la propagation des épidémies en Côte d'Ivoire\n"
        "pour sauver des vies et optimiser les ressources sanitaires?",
        fill_color=RGBColor(255, 220, 180),  # Orange léger
        line_color=ORANGE_CI,
        font_size=20,
        font_color=NOIR
    )
    problem_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Footer
    footer = slide.shapes.add_textbox(
        Inches(2), Inches(6.5), Inches(6), Inches(0.4)
    )
    p = footer.text_frame.paragraphs[0]
    p.text = "Orange Think Tank Challenge 2025 | Thème Santé"
    p.font.size = Pt(14)
    p.font.color.rgb = NOIR
    p.alignment = PP_ALIGN.CENTER

    # Numéro de page
    add_page_number(slide, 1)

    return slide


def create_slide_2_solution(prs):
    """
    SLIDE 2: La Solution - Architecture du Système
    """
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Titre
    add_title(slide, "Notre Solution: Un Système Prédictif Intelligent", 40)

    # Image architecture (centrée)
    if os.path.exists(ARCHITECTURE_IMAGE):
        left = Inches(1.5)
        top = Inches(1.5)
        pic = slide.shapes.add_picture(
            ARCHITECTURE_IMAGE, left, top, height=Inches(3.5)
        )

    # 3 Points clés (en bas)
    points_top = Inches(5.2)
    point_width = Inches(2.8)
    point_height = Inches(1.2)
    spacing = Inches(0.2)

    # Point 1: 15M d'abonnés
    add_shape_with_text(
        slide,
        Inches(0.5), points_top, point_width, point_height,
        "15M d'abonnés\n\n55% du marché ivoirien\nCouverture 95% du territoire",
        fill_color=GRIS_CLAIR,
        line_color=ORANGE_CI,
        font_size=14,
        font_color=NOIR
    )

    # Point 2: Prédictions
    add_shape_with_text(
        slide,
        Inches(0.5) + point_width + spacing, points_top, point_width, point_height,
        "Prédictions 7-14 jours\n\nAnticipation pour allocation\noptimale des ressources",
        fill_color=GRIS_CLAIR,
        line_color=ORANGE_CI,
        font_size=14,
        font_color=NOIR
    )

    # Point 3: 30 zones
    add_shape_with_text(
        slide,
        Inches(0.5) + 2*(point_width + spacing), points_top, point_width, point_height,
        "30 zones couvertes\n\n13 communes d'Abidjan\n+ 17 villes majeures",
        fill_color=GRIS_CLAIR,
        line_color=ORANGE_CI,
        font_size=14,
        font_color=NOIR
    )

    # Numéro de page
    add_page_number(slide, 2)

    return slide


def create_slide_3_anonymisation(prs):
    """
    SLIDE 3: Protection des Données - Anonymisation
    """
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Titre
    add_title(slide, "Confidentialité Absolue: Processus d'Anonymisation", 38)

    # Image anonymisation (centrée et grande)
    if os.path.exists(ANONYMISATION_IMAGE):
        left = Inches(1.5)
        top = Inches(1.3)
        pic = slide.shapes.add_picture(
            ANONYMISATION_IMAGE, left, top, height=Inches(4.5)
        )

    # Encadré garanties (en bas)
    guarantees_box = add_shape_with_text(
        slide,
        Inches(0.8), Inches(6), Inches(8.4), Inches(0.9),
        "✓ K-anonymat k≥50 personnes par groupe  •  ✓ Aucun identifiant personnel conservé\n"
        "✓ Conformité RGPD et législation ivoirienne  •  ✓ Comité d'éthique de supervision",
        fill_color=RGBColor(250, 200, 150),  # Orange très léger
        line_color=ORANGE_CI,
        font_size=14,
        font_color=NOIR
    )
    guarantees_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Numéro de page
    add_page_number(slide, 3)

    return slide


def create_slide_4_methodologie(prs):
    """
    SLIDE 4: Méthodologie - Comment ça Marche?
    """
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Titre
    add_title(slide, "Méthodologie: De la Donnée à la Prédiction", 40)

    # 4 étapes horizontales avec flèches
    box_width = Inches(2.0)
    box_height = Inches(3.8)
    box_top = Inches(1.8)
    spacing = Inches(0.15)
    arrow_width = Inches(0.3)

    # Étape 1: COLLECTE
    box1 = add_shape_with_text(
        slide,
        Inches(0.5), box_top, box_width, box_height,
        "ÉTAPE 1: COLLECTE\n\n"
        "• 15M d'abonnés Orange CI\n"
        "• Données CDR (Call Detail Records)\n"
        "• Matrices origine-destination quotidiennes",
        fill_color=GRIS_CLAIR,
        line_color=ORANGE_CI,
        font_size=13,
        font_color=NOIR
    )

    # Flèche 1
    add_text_box(
        slide,
        Inches(0.5) + box_width, box_top + box_height/2 - Inches(0.2),
        arrow_width, Inches(0.4),
        "→",
        font_size=32,
        color=ORANGE_CI,
        alignment=PP_ALIGN.CENTER
    )

    # Étape 2: TRAITEMENT
    x2 = Inches(0.5) + box_width + arrow_width + spacing
    box2 = add_shape_with_text(
        slide,
        x2, box_top, box_width, box_height,
        "ÉTAPE 2: TRAITEMENT\n\n"
        "• Anonymisation K-anonymat k≥50\n"
        "• Agrégation spatiale (30 zones)\n"
        "• Agrégation temporelle",
        fill_color=GRIS_CLAIR,
        line_color=ORANGE_CI,
        font_size=13,
        font_color=NOIR
    )

    # Flèche 2
    add_text_box(
        slide,
        x2 + box_width, box_top + box_height/2 - Inches(0.2),
        arrow_width, Inches(0.4),
        "→",
        font_size=32,
        color=ORANGE_CI,
        alignment=PP_ALIGN.CENTER
    )

    # Étape 3: MODÉLISATION
    x3 = x2 + box_width + arrow_width + spacing
    box3 = add_shape_with_text(
        slide,
        x3, box_top, box_width, box_height,
        "ÉTAPE 3: MODÉLISATION\n\n"
        "• Modèle SEIR métapopulationnel\n"
        "• Intégration flux de mobilité\n"
        "• Calibration sur épidémies passées",
        fill_color=GRIS_CLAIR,
        line_color=ORANGE_CI,
        font_size=13,
        font_color=NOIR
    )

    # Flèche 3
    add_text_box(
        slide,
        x3 + box_width, box_top + box_height/2 - Inches(0.2),
        arrow_width, Inches(0.4),
        "→",
        font_size=32,
        color=ORANGE_CI,
        alignment=PP_ALIGN.CENTER
    )

    # Étape 4: PRÉDICTIONS
    x4 = x3 + box_width + arrow_width + spacing
    box4 = add_shape_with_text(
        slide,
        x4, box_top, box_width, box_height,
        "ÉTAPE 4: PRÉDICTIONS\n\n"
        "• Prédictions J+7 et J+14\n"
        "• Précision visée >75%\n"
        "• Identification zones à risque",
        fill_color=GRIS_CLAIR,
        line_color=ORANGE_CI,
        font_size=13,
        font_color=NOIR
    )

    # Numéro de page
    add_page_number(slide, 4)

    return slide


def create_slide_5_livrables(prs):
    """
    SLIDE 5: Livrables & Impacts
    """
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Titre
    add_title(slide, "Livrables & Valeur Créée", 44)

    # Séparateur vertical orange
    separator = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(5), Inches(1.5), Pt(3), Inches(5)
    )
    separator.fill.solid()
    separator.fill.fore_color.rgb = ORANGE_CI
    separator.line.color.rgb = ORANGE_CI

    # COLONNE 1: LIVRABLES (gauche)
    col1_left = Inches(0.3)
    col1_width = Inches(4.5)

    # Titre colonne 1
    add_text_box(
        slide,
        col1_left, Inches(1.3), col1_width, Inches(0.4),
        "LIVRABLES CONCRETS",
        font_size=22,
        bold=True,
        color=ORANGE_CI,
        alignment=PP_ALIGN.CENTER
    )

    # Livrable 1
    add_shape_with_text(
        slide,
        col1_left, Inches(1.9), col1_width, Inches(1.0),
        "1. Dashboard Interactif Web\n"
        "Carte de Côte d'Ivoire avec zones colorées par risque, "
        "flux de mobilité animés, prédictions J+7/J+14",
        fill_color=BLANC,
        line_color=ORANGE_CI,
        font_size=12,
        font_color=NOIR
    )

    # Livrable 2
    add_shape_with_text(
        slide,
        col1_left, Inches(3.1), col1_width, Inches(0.7),
        "2. API d'Intégration\n"
        "Connexion aux systèmes du Ministère de la Santé, "
        "mises à jour quotidiennes",
        fill_color=BLANC,
        line_color=ORANGE_CI,
        font_size=12,
        font_color=NOIR
    )

    # Livrable 3
    add_shape_with_text(
        slide,
        col1_left, Inches(4.0), col1_width, Inches(0.7),
        "3. Système d'Alertes\n"
        "Notifications automatiques (email, SMS), "
        "détection zones critiques (score ≥85)",
        fill_color=BLANC,
        line_color=ORANGE_CI,
        font_size=12,
        font_color=NOIR
    )

    # Livrable 4
    add_shape_with_text(
        slide,
        col1_left, Inches(4.9), col1_width, Inches(0.6),
        "4. Formation & Documentation\n"
        "Formation 20-30 professionnels de santé, guides d'utilisation",
        fill_color=BLANC,
        line_color=ORANGE_CI,
        font_size=12,
        font_color=NOIR
    )

    # COLONNE 2: IMPACTS (droite)
    col2_left = Inches(5.2)
    col2_width = Inches(4.5)

    # Titre colonne 2
    add_text_box(
        slide,
        col2_left, Inches(1.3), col2_width, Inches(0.4),
        "IMPACTS MESURABLES",
        font_size=22,
        bold=True,
        color=ORANGE_CI,
        alignment=PP_ALIGN.CENTER
    )

    # Section Population
    add_text_box(
        slide,
        col2_left, Inches(1.85), col2_width, Inches(0.35),
        "Pour la Population Ivoirienne:",
        font_size=14,
        bold=True,
        color=NOIR
    )

    pop_box = slide.shapes.add_textbox(
        col2_left, Inches(2.25), col2_width, Inches(1.8)
    )
    tf = pop_box.text_frame
    tf.word_wrap = True

    # Points population
    for text in [
        "Réduction 30% du temps de réponse aux épidémies",
        "Allocation optimisée des ressources sanitaires",
        "Sauvetage de vies grâce à l'anticipation",
        "Restrictions ciblées vs généralisées"
    ]:
        p = tf.add_paragraph() if tf.paragraphs[0].text else tf.paragraphs[0]
        p.text = "• " + text
        p.font.size = Pt(13)
        p.font.color.rgb = NOIR
        p.space_after = Pt(4)

    # Section Orange CI
    add_text_box(
        slide,
        col2_left, Inches(4.2), col2_width, Inches(0.35),
        "Pour Orange Côte d'Ivoire:",
        font_size=14,
        bold=True,
        color=NOIR
    )

    orange_box = slide.shapes.add_textbox(
        col2_left, Inches(4.6), col2_width, Inches(1.5)
    )
    tf = orange_box.text_frame
    tf.word_wrap = True

    # Points Orange CI
    for text in [
        "Leadership RSE et innovation sociale",
        "Valorisation des actifs data existants",
        "Partenariats institutionnels stratégiques",
        "Modèle exportable (Orange Africa)"
    ]:
        p = tf.add_paragraph() if tf.paragraphs[0].text else tf.paragraphs[0]
        p.text = "• " + text
        p.font.size = Pt(13)
        p.font.color.rgb = NOIR
        p.space_after = Pt(4)

    # Numéro de page
    add_page_number(slide, 5)

    return slide


def create_slide_6_risques(prs):
    """
    SLIDE 6: Risques & Stratégies de Mitigation
    """
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Titre
    add_title(slide, "Risques Identifiés & Solutions", 44)

    # Tableau 2 colonnes (4 lignes de risques)
    col1_left = Inches(0.5)
    col2_left = Inches(5.2)
    col_width = Inches(4.5)
    row_height = Inches(1.15)
    top_start = Inches(1.5)
    row_spacing = Inches(0.05)

    risques = [
        {
            "risque": "1. Vie privée des citoyens",
            "mitigation": "• K-anonymat renforcé (k≥50)\n"
                         "• Suppression totale des identifiants\n"
                         "• Comité d'éthique indépendant\n"
                         "• Audits annuels par tiers"
        },
        {
            "risque": "2. Précision des modèles",
            "mitigation": "• Calibration sur épidémies passées (COVID, Dengue)\n"
                         "• Validation scientifique continue\n"
                         "• Intervalle de confiance ±15%\n"
                         "• Amélioration continue"
        },
        {
            "risque": "3. Adoption par utilisateurs",
            "mitigation": "• Formation complète (20-30 personnes)\n"
                         "• Interface intuitive et visuelle\n"
                         "• Support technique dédié\n"
                         "• Co-conception avec Ministère Santé"
        },
        {
            "risque": "4. Pérennité du système",
            "mitigation": "• Transfert de compétences progressif\n"
                         "• Documentation exhaustive\n"
                         "• Modèle économique long terme\n"
                         "• Infrastructure scalable (cloud)"
        }
    ]

    for i, risk_item in enumerate(risques):
        row_top = top_start + i * (row_height + row_spacing)

        # Colonne RISQUE
        risk_box = add_shape_with_text(
            slide,
            col1_left, row_top, col_width, row_height,
            risk_item["risque"],
            fill_color=RGBColor(255, 240, 230),  # Orange très léger
            line_color=ORANGE_CI,
            font_size=14,
            font_color=NOIR
        )
        risk_box.text_frame.paragraphs[0].font.bold = True

        # Colonne MITIGATION
        mitigation_box = add_shape_with_text(
            slide,
            col2_left, row_top, col_width, row_height,
            risk_item["mitigation"],
            fill_color=BLANC,
            line_color=ORANGE_CI,
            font_size=12,
            font_color=NOIR
        )

    # Numéro de page
    add_page_number(slide, 6)

    return slide


def create_slide_7_prototype(prs):
    """
    SLIDE 7: Démonstration du Prototype
    """
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Titre
    add_title(slide, "Prototype Fonctionnel", 44)

    # Message pour capture d'écran
    placeholder_box = add_shape_with_text(
        slide,
        Inches(1.5), Inches(2), Inches(7), Inches(3.5),
        "CAPTURE D'ÉCRAN DU DASHBOARD\n\n"
        "À insérer ici:\n"
        "• Carte 3D de Côte d'Ivoire avec zones colorées\n"
        "• 4 KPI principaux (Cas actifs, Mobilité, Zones à risque, Prédiction J+7)\n"
        "• Graphique d'évolution temporelle\n"
        "• Tableau des régions\n\n"
        "(Veuillez capturer une vue du dashboard et remplacer cette zone)",
        fill_color=GRIS_CLAIR,
        line_color=ORANGE_CI,
        font_size=16,
        font_color=NOIR
    )
    placeholder_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Call to action
    cta_box = add_shape_with_text(
        slide,
        Inches(2.5), Inches(6), Inches(5), Inches(0.6),
        "Démonstration interactive du système disponible",
        fill_color=ORANGE_CI,
        line_color=ORANGE_CI,
        font_size=20,
        font_color=BLANC
    )
    cta_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    cta_box.text_frame.paragraphs[0].font.bold = True

    # Numéro de page
    add_page_number(slide, 7)

    return slide


# =============================================================================
# FONCTION PRINCIPALE
# =============================================================================

def create_presentation():
    """
    Fonction principale qui crée la présentation complète
    """
    print("🚀 Début de la création de la présentation...")

    # Créer une nouvelle présentation
    prs = Presentation()
    prs.slide_width = Inches(10)  # 16:9
    prs.slide_height = Inches(7.5)

    print("📄 Création de la Slide 1: Titre et Problématique...")
    create_slide_1_title(prs)

    print("📄 Création de la Slide 2: Solution et Architecture...")
    create_slide_2_solution(prs)

    print("📄 Création de la Slide 3: Anonymisation...")
    create_slide_3_anonymisation(prs)

    print("📄 Création de la Slide 4: Méthodologie...")
    create_slide_4_methodologie(prs)

    print("📄 Création de la Slide 5: Livrables et Impacts...")
    create_slide_5_livrables(prs)

    print("📄 Création de la Slide 6: Risques et Mitigations...")
    create_slide_6_risques(prs)

    print("📄 Création de la Slide 7: Prototype...")
    create_slide_7_prototype(prs)

    # Sauvegarder la présentation
    output_filename = "Orange_Think_Tank_2025_Prediction_Epidemies.pptx"
    prs.save(output_filename)

    print(f"\n✅ Présentation créée avec succès: {output_filename}")
    print(f"📊 Nombre de slides: {len(prs.slides)}")
    print("\n📋 Résumé:")
    print("  - Slide 1: Titre et Problématique")
    print("  - Slide 2: Solution et Architecture (avec architecture.png)")
    print("  - Slide 3: Anonymisation (avec anonymisation.png)")
    print("  - Slide 4: Méthodologie en 4 étapes")
    print("  - Slide 5: Livrables et Impacts (2 colonnes)")
    print("  - Slide 6: Risques et Mitigations")
    print("  - Slide 7: Prototype")
    print("\n⚠️  N'oubliez pas:")
    print("  - Vérifier que architecture.png et anonymisation.png sont dans le même dossier")
    print("  - Ajouter une capture d'écran du dashboard sur la slide 7")
    print("  - Vérifier les couleurs et ajuster si nécessaire")

    return output_filename


if __name__ == "__main__":
    try:
        create_presentation()
    except FileNotFoundError as e:
        print(f"\n❌ Erreur: Fichier image manquant - {e}")
        print("⚠️  Assurez-vous que les fichiers suivants sont présents:")
        print("   - architecture.png")
        print("   - anonymisation.png")
    except Exception as e:
        print(f"\n❌ Erreur lors de la création: {e}")
        import traceback
        traceback.print_exc()
